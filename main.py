import os
import io
import secrets
from datetime import datetime, timedelta

from fastapi import FastAPI, UploadFile, File, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from PIL import Image

# =========================
# Assets (solo logo)
# =========================
ASSETS_DIR = "assets"
DEFAULT_LOGO_PATH = os.path.join(ASSETS_DIR, "logo.png")

os.makedirs(ASSETS_DIR, exist_ok=True)

# Generar logo por defecto si no existe
if not os.path.exists(DEFAULT_LOGO_PATH):
    img = Image.new("RGB", (600, 600), (255, 255, 255))
    try:
        from PIL import ImageDraw
        draw = ImageDraw.Draw(img)
        draw.ellipse((100, 100, 500, 500), fill=(133, 78, 197))  # círculo morado
    except Exception:
        pass
    img.save(DEFAULT_LOGO_PATH, format="PNG")

# =========================
# Configuración de formateo
# =========================
TITLE_SIZE = 20
TITLE_COLOR = RGBColor(0, 0, 0)   # negro
NORMAL_COLOR = RGBColor(0, 0, 0)     # negro
FONT = "Century Gothic"
EMU_PER_INCH = 914400

# =========================
# FastAPI
# =========================
app = FastAPI(title="Formateador de Presentaciones PPTX")

ALLOWED_ORIGINS = [
    "https://www.dipli.ai",
    "https://dipli.ai",
    "https://isagarcivill09.wixsite.com/turop",
    "https://isagarcivill09.wixsite.com/turop/tienda",
    "https://isagarcivill09-wixsite-com.filesusr.com",
    "https://www.dipli.ai/preparaci%C3%B3n",
    "https://www-dipli-ai.filesusr.com",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    allow_origin_regex=r"https://.*\.filesusr\.com",
)

# =========================
# Descargas temporales
# =========================
DOWNLOADS: dict[str, tuple[bytes, str, str, datetime]] = {}
DOWNLOAD_TTL_SECS = 900
PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def cleanup_downloads():
    now = datetime.utcnow()
    expired = [t for t, (_, _, _, exp) in DOWNLOADS.items() if exp <= now]
    for t in expired:
        DOWNLOADS.pop(t, None)


def register_download(data: bytes, filename: str, media_type: str) -> str:
    cleanup_downloads()
    token = secrets.token_urlsafe(16)
    expires_at = datetime.utcnow() + timedelta(seconds=DOWNLOAD_TTL_SECS)
    DOWNLOADS[token] = (data, filename, media_type, expires_at)
    return token

# =========================
# Formato de texto
# =========================
def format_run(run):
    run.font.name = FONT
    if run.font.size and run.font.size.pt >= TITLE_SIZE:
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR
    else:
        run.font.bold = False
        run.font.color.rgb = NORMAL_COLOR


def apply_rules(shape):
    if hasattr(shape, "text_frame") and shape.text_frame:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                format_run(r)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        format_run(r)
    if getattr(shape, "shape_type", None) == 6:
        for s in shape.shapes:
            apply_rules(s)

# =========================
# Procesamiento principal
# =========================
def process_presentation(file_bytes: bytes, filename: str) -> bytes:
    prs = Presentation(io.BytesIO(file_bytes))

    if len(prs.slides) > 0:
        slide0 = prs.slides[0]

        # Buscar shape principal con texto (título)
        title_shape = None
        for s in slide0.shapes:
            if hasattr(s, "text_frame") and s.text_frame and s.text_frame.text.strip():
                title_shape = s
                break

        if title_shape:
            # Calcular zona real de texto (solo las líneas ocupadas)
            text_height = 0
            line_spacing = 0
            for p in title_shape.text_frame.paragraphs:
                if p.text.strip():
                    font_size = None
                    for r in p.runs:
                        if r.font.size:
                            font_size = r.font.size.pt
                            break
                    if font_size:
                        text_height += font_size * 1.3
            if text_height == 0:
                text_height = (title_shape.height / EMU_PER_INCH) * 0.6

            # Coordenadas del shape
            ref_left = title_shape.left
            ref_top = title_shape.top
            ref_w = title_shape.width
            ref_h = title_shape.height

            # Insertar logo justo debajo del texto visible
            logo_w_in = max(1.5, min((ref_w / EMU_PER_INCH) * 0.4, 5))
            pic = slide0.shapes.add_picture(
                DEFAULT_LOGO_PATH,
                left=Inches(0), top=Inches(0), width=Inches(logo_w_in)
            )
            img_w = pic.width
            img_h = pic.height

            # Centrar horizontalmente con el texto y ubicar justo debajo
            pic.left = ref_left + (ref_w - img_w) // 2
            pic.top = ref_top + int(text_height * 1.1 * EMU_PER_INCH / 72)

    # Aplicar reglas a todas las diapositivas
    for slide in prs.slides:
        for shp in slide.shapes:
            try:
                apply_rules(shp)
            except Exception:
                pass

    # Aplicar a layouts
    for layout in prs.slide_layouts:
        for shp in layout.shapes:
            try:
                apply_rules(shp)
            except Exception:
                pass

    out_bytes = io.BytesIO()
    prs.save(out_bytes)
    out_bytes.seek(0)
    return out_bytes.getvalue()

# =========================
# Endpoints
# =========================
@app.post("/procesar/")
async def procesar_pptx(request: Request, file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="El archivo debe ser un .pptx válido")

    file_bytes = await file.read()
    try:
        result_bytes = process_presentation(file_bytes, file.filename)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error procesando PPTX: {e}")

    final_name = file.filename.replace(".pptx", "_FORMATEADO.pptx")
    token = register_download(result_bytes, final_name, PPTX_MEDIA_TYPE)

    base_url = str(request.base_url).rstrip('/')
    download_url = f"{base_url}/download/{token}"
    return {"download_url": download_url, "expires_in_seconds": DOWNLOAD_TTL_SECS}


@app.get("/download/{token}")
def download_token(token: str):
    cleanup_downloads()
    item = DOWNLOADS.get(token)
    if not item:
        raise HTTPException(status_code=404, detail="Link expirado o inválido")
    data, filename, media_type, exp = item
    if exp <= datetime.utcnow():
        DOWNLOADS.pop(token, None)
        raise HTTPException(status_code=410, detail="Link expirado")

    headers = {
        "Content-Disposition": f'attachment; filename="{filename}"',
        "Cache-Control": "no-store",
    }
    return StreamingResponse(io.BytesIO(data), media_type=media_type, headers=headers)


@app.get("/")
async def root():
    return {"message": "API de Formateo de PPTX funcionando", "version": "1.0.1"}


@app.get("/health")
async def health_check():
    return {"status": "healthy", "message": "API funcionando correctamente"}


